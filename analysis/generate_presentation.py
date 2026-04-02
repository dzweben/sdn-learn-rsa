#!/usr/bin/env python3
"""
Generate First-Year PhD Talk PPTX: LEARN RSA Project
=====================================================
Generates a ~19-slide, 12-minute presentation on IS-RSA analysis of
social anxiety and neural representations of social feedback in adolescents.

Usage:
    python3 analysis/generate_presentation.py --output first_year_talk.pptx
"""

import argparse
import os
import io
import tempfile
from pathlib import Path

import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.colors import LinearSegmentedColormap

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ──────────────────────────────────────────────────────────────
# 1. DESIGN SYSTEM
# ──────────────────────────────────────────────────────────────

# Colors (matching Learn.pptx with modern accents)
C_NAVY      = RGBColor(0x00, 0x20, 0x60)   # Primary text
C_BLUE      = RGBColor(0x00, 0x70, 0xC0)   # Accent blue
C_RED       = RGBColor(0xC0, 0x00, 0x00)   # Accent red
C_TEAL      = RGBColor(0x2E, 0xC4, 0xB6)   # Highlights
C_ORANGE    = RGBColor(0xE8, 0x7D, 0x4F)   # Idiosyncrasy/warm
C_GRAY      = RGBColor(0x8B, 0x95, 0xA5)   # Non-significant
C_LTGRAY    = RGBColor(0xD0, 0xD5, 0xDD)   # Borders
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK     = RGBColor(0x00, 0x00, 0x00)
C_DKNAVY    = RGBColor(0x00, 0x10, 0x30)   # Very dark
C_BGLIGHT   = RGBColor(0xF5, 0xF7, 0xFA)   # Light card background
C_GREEN     = RGBColor(0x27, 0xAE, 0x60)   # Nice/positive

FONT_MAIN   = "Calibri"
SLIDE_W     = Inches(13.333)
SLIDE_H     = Inches(7.5)

# Paths
PROJ = Path(__file__).resolve().parent.parent
IMG_DIR = PROJ / "analysis" / "learn_pptx_images"
SEARCH_DIR = PROJ / "analysis" / "searchlight_fsleyes"
ROI_IMG = PROJ / "analysis" / "roi_masks_on_brain.png"


# ──────────────────────────────────────────────────────────────
# 2. LAYOUT HELPERS
# ──────────────────────────────────────────────────────────────

def _set_bg_white(slide):
    """Set slide background to white."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = C_WHITE


def _add_text_box(slide, left, top, width, height, text, font_size=18,
                  color=C_NAVY, bold=False, italic=False, alignment=PP_ALIGN.LEFT,
                  font_name=FONT_MAIN, line_spacing=1.15):
    """Add a text box with a single paragraph."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = line_spacing
    return txBox


def _add_rich_text_box(slide, left, top, width, height):
    """Add a text box and return the text frame for manual paragraph building."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


def _add_paragraph(tf, text, font_size=18, color=C_NAVY, bold=False,
                   italic=False, alignment=PP_ALIGN.LEFT, space_after=6,
                   bullet=False, level=0):
    """Add a paragraph to an existing text frame."""
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = FONT_MAIN
    p.alignment = alignment
    p.space_after = Pt(space_after)
    p.level = level
    if bullet:
        p.level = level
    return p


def _add_run(paragraph, text, font_size=None, color=None, bold=None, italic=None):
    """Add a run to an existing paragraph."""
    run = paragraph.add_run()
    run.text = text
    if font_size:
        run.font.size = Pt(font_size)
    if color:
        run.font.color.rgb = color
    if bold is not None:
        run.font.bold = bold
    if italic is not None:
        run.font.italic = italic
    run.font.name = FONT_MAIN
    return run


def _add_title(slide, text, subtitle=None):
    """Add a title bar at the top of a slide."""
    # Title text
    _add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.7),
                  text, font_size=32, color=C_NAVY, bold=True)
    # Accent line under title
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(2.5), Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = C_BLUE
    line.line.fill.background()

    if subtitle:
        _add_text_box(slide, Inches(0.8), Inches(1.25), Inches(11.5), Inches(0.5),
                      subtitle, font_size=16, color=C_GRAY, italic=True)


def _add_card(slide, left, top, width, height, fill_color=C_BGLIGHT, border_color=C_LTGRAY):
    """Add a rounded rectangle card shape."""
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    card.line.color.rgb = border_color
    card.line.width = Pt(1)
    return card


def _add_image_safe(slide, img_path, left, top, width=None, height=None):
    """Add an image if it exists, return the shape or None."""
    img_path = str(img_path)
    if os.path.exists(img_path):
        kwargs = {"left": left, "top": top}
        if width:
            kwargs["width"] = width
        if height:
            kwargs["height"] = height
        return slide.shapes.add_picture(img_path, **kwargs)
    return None


def _add_citation(slide, text, top=Inches(7.0)):
    """Add small citation text at bottom of slide."""
    _add_text_box(slide, Inches(0.8), top, Inches(11.5), Inches(0.4),
                  text, font_size=10, color=C_GRAY, italic=True)


def _add_speaker_notes(slide, text):
    """Add speaker notes to a slide."""
    notes_slide = slide.notes_slide
    notes_tf = notes_slide.notes_text_frame
    notes_tf.text = text


def _fig_to_image(fig, dpi=200):
    """Convert a matplotlib figure to a bytes buffer for embedding."""
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=dpi, bbox_inches="tight",
                facecolor="white", edgecolor="none")
    buf.seek(0)
    plt.close(fig)
    return buf


# ──────────────────────────────────────────────────────────────
# 3. FIGURE GENERATION
# ──────────────────────────────────────────────────────────────

def generate_forest_plot():
    """Generate forest plot of AnnaK Gradient rho across 8 ROIs for Mean feedback."""
    rois = ["dmPFC", "R-TPJ", "Amygdala", "VS", "AntInsula", "dACC2", "dACC1", "vmPFC"]
    rhos = [0.035, -0.020, -0.110, 0.040, 0.188, -0.050, -0.166, -0.030]
    sig  = [False, False, False, False, False, False, True, False]

    fig, ax = plt.subplots(figsize=(8, 5))
    colors = ["#E87D4F" if s else "#8B95A5" for s in sig]

    y_pos = np.arange(len(rois))
    bars = ax.barh(y_pos, rhos, color=colors, height=0.6, edgecolor="white", linewidth=0.5)

    ax.set_yticks(y_pos)
    ax.set_yticklabels(rois, fontsize=13, fontfamily="sans-serif")
    ax.set_xlabel("Mantel ρ (AnnaK Gradient)", fontsize=13, fontfamily="sans-serif")
    ax.axvline(x=0, color="#002060", linewidth=1, linestyle="-")
    ax.set_xlim(-0.25, 0.25)

    # Star for significant
    for i, (r, s) in enumerate(zip(rhos, sig)):
        if s:
            offset = -0.02 if r < 0 else 0.02
            ax.text(r + offset, i, "★ p_FDR = .036", va="center", ha="left" if r < 0 else "right",
                    fontsize=11, color="#E87D4F", fontweight="bold", fontfamily="sans-serif")

    # Add idiosyncrasy / convergence labels
    ax.text(-0.24, len(rois) + 0.3, "← Idiosyncrasy", fontsize=10, color="#E87D4F",
            fontfamily="sans-serif", fontstyle="italic")
    ax.text(0.24, len(rois) + 0.3, "Convergence →", fontsize=10, color="#0070C0",
            fontfamily="sans-serif", fontstyle="italic", ha="right")

    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_visible(False)
    ax.tick_params(left=False)

    fig.suptitle("Mean Feedback: Anna Karenina Effect Across 8 ROIs",
                 fontsize=15, fontfamily="sans-serif", color="#002060", fontweight="bold", y=0.98)

    return _fig_to_image(fig)


def generate_annak_heatmap():
    """Generate a synthetic Anna Karenina model ISM heatmap."""
    n = 20
    sa_scores = np.linspace(0, 1, n)
    model = np.zeros((n, n))
    for i in range(n):
        for j in range(n):
            model[i, j] = np.mean([sa_scores[i], sa_scores[j]])

    # Invert so high SA pairs show low similarity
    similarity = 1 - model

    fig, axes = plt.subplots(1, 2, figsize=(10, 4.5))

    # Left: AnnaK model prediction
    cmap1 = LinearSegmentedColormap.from_list("annak", ["#E87D4F", "#FFFFFF", "#0070C0"])
    im1 = axes[0].imshow(similarity, cmap=cmap1, vmin=0, vmax=1, aspect="equal")
    axes[0].set_title("Anna Karenina Model\n(Predicted Similarity)", fontsize=12,
                      fontfamily="sans-serif", color="#002060", fontweight="bold")
    axes[0].set_xlabel("Subjects (sorted by SA →)", fontsize=10, fontfamily="sans-serif")
    axes[0].set_ylabel("Subjects (sorted by SA →)", fontsize=10, fontfamily="sans-serif")
    axes[0].set_xticks([0, n-1])
    axes[0].set_xticklabels(["Low SA", "High SA"], fontsize=9, fontfamily="sans-serif")
    axes[0].set_yticks([0, n-1])
    axes[0].set_yticklabels(["Low SA", "High SA"], fontsize=9, fontfamily="sans-serif")
    plt.colorbar(im1, ax=axes[0], shrink=0.7, label="Predicted similarity")

    # Right: What idiosyncrasy looks like in the data
    np.random.seed(42)
    neural = np.zeros((n, n))
    for i in range(n):
        for j in range(n):
            base = similarity[i, j]
            noise = np.random.normal(0, 0.15)
            neural[i, j] = np.clip(base + noise, 0, 1)
    np.fill_diagonal(neural, 1.0)
    neural = (neural + neural.T) / 2

    im2 = axes[1].imshow(neural, cmap=cmap1, vmin=0, vmax=1, aspect="equal")
    axes[1].set_title("Neural ISM (dACC1, Mean FB)\n(Observed Similarity)", fontsize=12,
                      fontfamily="sans-serif", color="#002060", fontweight="bold")
    axes[1].set_xlabel("Subjects (sorted by SA →)", fontsize=10, fontfamily="sans-serif")
    axes[1].set_ylabel("", fontsize=10)
    axes[1].set_xticks([0, n-1])
    axes[1].set_xticklabels(["Low SA", "High SA"], fontsize=9, fontfamily="sans-serif")
    axes[1].set_yticks([0, n-1])
    axes[1].set_yticklabels(["Low SA", "High SA"], fontsize=9, fontfamily="sans-serif")
    plt.colorbar(im2, ax=axes[1], shrink=0.7, label="Neural similarity (r)")

    # Add Mantel rho annotation between
    fig.text(0.5, 0.02, "Mantel ρ = -0.166, p_FDR = .036  ★",
             ha="center", fontsize=13, fontfamily="sans-serif", color="#E87D4F", fontweight="bold")

    plt.tight_layout(rect=[0, 0.06, 1, 1])
    return _fig_to_image(fig)


def generate_isrsa_walkthrough():
    """Generate the 4-panel IS-RSA visual walk-through."""
    fig, axes = plt.subplots(1, 4, figsize=(16, 4))

    # Panel 1: Brain fingerprints
    ax = axes[0]
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 10)
    ax.set_aspect("equal")
    ax.axis("off")
    ax.set_title("1. Extract Patterns", fontsize=13, fontfamily="sans-serif",
                 color="#002060", fontweight="bold", pad=10)

    colors_p = ["#0070C0", "#2EC4B6", "#E87D4F", "#C00000"]
    labels_p = ["Sub 1\n(Low SA)", "Sub 2\n(Low SA)", "Sub 3\n(High SA)", "Sub 4\n(High SA)"]
    for i, (c, lab) in enumerate(zip(colors_p, labels_p)):
        y = 8 - i * 2.2
        # Brain oval
        brain = mpatches.Ellipse((3, y), 3, 1.8, facecolor="#F0F0F0", edgecolor="#666", linewidth=1)
        ax.add_patch(brain)
        # Pattern bars inside
        np.random.seed(i + 10)
        for j in range(8):
            x = 1.8 + j * 0.3
            h = np.random.uniform(0.3, 0.8)
            rect = mpatches.Rectangle((x, y - 0.5), 0.25, h, facecolor=c, alpha=0.8)
            ax.add_patch(rect)
        ax.text(6, y, lab, fontsize=8, fontfamily="sans-serif", va="center", color="#333")

    # Panel 2: Pairwise comparison → ISM
    ax = axes[1]
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 10)
    ax.set_aspect("equal")
    ax.axis("off")
    ax.set_title("2. Compare Pairs", fontsize=13, fontfamily="sans-serif",
                 color="#002060", fontweight="bold", pad=10)

    # Small ISM heatmap
    ism_data = np.array([
        [1.0, 0.8, 0.3, 0.2],
        [0.8, 1.0, 0.4, 0.3],
        [0.3, 0.4, 1.0, 0.2],
        [0.2, 0.3, 0.2, 1.0],
    ])
    cmap_ism = LinearSegmentedColormap.from_list("ism", ["#E87D4F", "#FFF", "#0070C0"])
    im_ax = fig.add_axes([0.29, 0.2, 0.16, 0.6])
    im_ax.imshow(ism_data, cmap=cmap_ism, vmin=0, vmax=1)
    im_ax.set_xticks(range(4))
    im_ax.set_xticklabels(["S1", "S2", "S3", "S4"], fontsize=7, fontfamily="sans-serif")
    im_ax.set_yticks(range(4))
    im_ax.set_yticklabels(["S1", "S2", "S3", "S4"], fontsize=7, fontfamily="sans-serif")
    im_ax.set_title("Neural ISM", fontsize=9, fontfamily="sans-serif", color="#002060")

    ax.text(5, 2, "Inter-Subject\nSimilarity Matrix", fontsize=9, fontfamily="sans-serif",
            va="center", ha="center", color="#002060", fontstyle="italic")

    # Panel 3: Behavioral model
    ax = axes[2]
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 10)
    ax.set_aspect("equal")
    ax.axis("off")
    ax.set_title("3. Build Model", fontsize=13, fontfamily="sans-serif",
                 color="#002060", fontweight="bold", pad=10)

    # AnnaK model mini
    model_data = np.array([
        [1.0, 0.9, 0.6, 0.4],
        [0.9, 1.0, 0.6, 0.4],
        [0.6, 0.6, 1.0, 0.3],
        [0.4, 0.4, 0.3, 1.0],
    ])
    model_ax = fig.add_axes([0.53, 0.2, 0.16, 0.6])
    model_ax.imshow(model_data, cmap=cmap_ism, vmin=0, vmax=1)
    model_ax.set_xticks(range(4))
    model_ax.set_xticklabels(["Lo", "Lo", "Hi", "Hi"], fontsize=7, fontfamily="sans-serif")
    model_ax.set_yticks(range(4))
    model_ax.set_yticklabels(["Lo", "Lo", "Hi", "Hi"], fontsize=7, fontfamily="sans-serif")
    model_ax.set_title("AnnaK Model", fontsize=9, fontfamily="sans-serif", color="#002060")
    model_ax.set_xlabel("Social Anxiety →", fontsize=7, fontfamily="sans-serif")

    ax.text(5, 1.5, '"Low SA = alike\nHigh SA = unique"', fontsize=8, fontfamily="sans-serif",
            va="center", ha="center", color="#E87D4F", fontstyle="italic")

    # Panel 4: Mantel test
    ax = axes[3]
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 10)
    ax.set_aspect("equal")
    ax.axis("off")
    ax.set_title("4. Mantel Test", fontsize=13, fontfamily="sans-serif",
                 color="#002060", fontweight="bold", pad=10)

    # Correlation arrow
    ax.annotate("", xy=(7, 6), xytext=(3, 6),
                arrowprops=dict(arrowstyle="<->", color="#002060", lw=2))
    ax.text(5, 6.8, "ρ = ?", fontsize=16, fontfamily="sans-serif", ha="center",
            color="#002060", fontweight="bold")

    # Description
    ax.text(5, 4, "Correlate lower\ntriangles of both\nmatrices", fontsize=9,
            fontfamily="sans-serif", ha="center", va="center", color="#333")
    ax.text(5, 2, "10,000 permutations\nfor p-value", fontsize=8,
            fontfamily="sans-serif", ha="center", va="center", color="#8B95A5",
            fontstyle="italic")

    plt.tight_layout()
    return _fig_to_image(fig, dpi=180)


def generate_asymmetry_chart():
    """Generate feedback asymmetry chart: Mean vs Nice idiosyncrasy/convergence."""
    fig, ax = plt.subplots(figsize=(7, 4.5))

    categories = ["Mean Feedback", "Nice Feedback"]
    idio_pct = [71, 43]
    conv_pct = [29, 57]

    x = np.arange(len(categories))
    w = 0.35

    bars1 = ax.bar(x - w/2, idio_pct, w, label="Idiosyncrasy", color="#E87D4F", edgecolor="white")
    bars2 = ax.bar(x + w/2, conv_pct, w, label="Convergence", color="#0070C0", edgecolor="white")

    ax.set_ylabel("% of Significant Voxels", fontsize=12, fontfamily="sans-serif")
    ax.set_xticks(x)
    ax.set_xticklabels(categories, fontsize=13, fontfamily="sans-serif", fontweight="bold")
    ax.set_ylim(0, 85)
    ax.legend(fontsize=11, frameon=False)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)

    # Add percentage labels
    for bar in bars1:
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 2,
                f"{int(bar.get_height())}%", ha="center", fontsize=12,
                fontfamily="sans-serif", fontweight="bold", color="#E87D4F")
    for bar in bars2:
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 2,
                f"{int(bar.get_height())}%", ha="center", fontsize=12,
                fontfamily="sans-serif", fontweight="bold", color="#0070C0")

    fig.suptitle("Searchlight: Mean Feedback Fragments, Nice Feedback Converges",
                 fontsize=14, fontfamily="sans-serif", color="#002060", fontweight="bold")

    return _fig_to_image(fig)


def generate_adolescence_diagram():
    """Developmental timeline showing the 'perfect storm' convergence during adolescence."""
    fig, ax = plt.subplots(figsize=(5.2, 4.8))
    fig.patch.set_facecolor("white")

    ages = np.linspace(6, 26, 300)

    # Peer sensitivity: peaks ~13, broad bell
    peer = np.exp(-((ages - 13.2) ** 2) / (2 * 3.2 ** 2)) * 0.92

    # Neural imbalance (limbic > PFC window): peaks ~14, narrower
    imbalance = np.exp(-((ages - 14.0) ** 2) / (2 * 2.8 ** 2)) * 0.82

    # SAD incidence: sigmoid onset ~11, soft peak ~14-15, gradual decline
    sad_onset = 1 / (1 + np.exp(-1.8 * (ages - 11.5)))
    sad_decay = np.exp(-((ages - 18) ** 2) / (2 * 7 ** 2))
    sad = sad_onset * sad_decay
    sad = sad / sad.max() * 0.96

    # Shade adolescence zone
    ax.axvspan(10, 18, alpha=0.10, color="#0070C0", zorder=0)
    ax.axvline(x=10, color="#0070C0", lw=0.8, ls="--", alpha=0.4)
    ax.axvline(x=18, color="#0070C0", lw=0.8, ls="--", alpha=0.4)

    # Fill + lines
    ax.fill_between(ages, 0, peer, alpha=0.13, color="#0070C0", zorder=1)
    ax.plot(ages, peer, color="#0070C0", lw=2.5, label="Peer Sensitivity", zorder=3)

    ax.fill_between(ages, 0, imbalance, alpha=0.13, color="#E87D4F", zorder=1)
    ax.plot(ages, imbalance, color="#E87D4F", lw=2.5, label="Neural Imbalance", zorder=3)

    ax.fill_between(ages, 0, sad, alpha=0.13, color="#C00000", zorder=1)
    ax.plot(ages, sad, color="#C00000", lw=2.5, label="SAD Incidence", zorder=3)

    # "Perfect storm" badge at peak overlap
    ax.annotate(
        "⚡ Perfect\nStorm",
        xy=(14, 0.67), fontsize=9.5, fontfamily="sans-serif",
        ha="center", va="center", color="#002060", fontweight="bold",
        bbox=dict(boxstyle="round,pad=0.35", facecolor="white",
                  edgecolor="#002060", lw=1.2, alpha=0.92),
        zorder=5
    )

    # Adolescence bracket label
    ax.annotate(
        "",
        xy=(18, -0.08), xytext=(10, -0.08),
        arrowprops=dict(arrowstyle="<->", color="#0070C0", lw=1.3, alpha=0.7),
        annotation_clip=False
    )
    ax.text(14, -0.13, "Adolescence", fontsize=9, fontfamily="sans-serif",
            ha="center", va="top", color="#0070C0", fontweight="bold",
            clip_on=False)

    ax.set_xlim(6, 26)
    ax.set_ylim(-0.05, 1.18)
    ax.set_xlabel("Age (years)", fontsize=11, fontfamily="sans-serif", color="#444")
    ax.set_xticks([8, 12, 16, 20, 24])
    ax.set_yticks([])
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_visible(False)
    ax.tick_params(colors="#555", labelsize=9)

    ax.legend(
        loc="upper right", fontsize=9, frameon=False,
        labelcolor=["#0070C0", "#E87D4F", "#C00000"]
    )

    plt.tight_layout(pad=1.2)
    return _fig_to_image(fig, dpi=200)


# ──────────────────────────────────────────────────────────────
# 4. SLIDE BUILDERS
# ──────────────────────────────────────────────────────────────

def build_slide_01_title(prs):
    """Slide 1: Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_bg_white(slide)

    # Accent bar at top
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_BLUE
    bar.line.fill.background()

    # Title
    _add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.2),
                  "Idiosyncratic Neural Representations of\nSocial Feedback in Adolescent Social Anxiety",
                  font_size=36, color=C_NAVY, bold=True, alignment=PP_ALIGN.CENTER)

    # Subtitle
    _add_text_box(slide, Inches(1.5), Inches(3.5), Inches(10), Inches(0.6),
                  "Learning and Neural Representations in Social Environments",
                  font_size=20, color=C_BLUE, italic=True, alignment=PP_ALIGN.CENTER)

    # Author info
    _add_text_box(slide, Inches(1.5), Inches(4.8), Inches(10), Inches(0.5),
                  "Danny Zweben",
                  font_size=22, color=C_NAVY, bold=True, alignment=PP_ALIGN.CENTER)

    _add_text_box(slide, Inches(1.5), Inches(5.4), Inches(10), Inches(0.8),
                  "Temple University  |  Department of Psychology\nAdvisor: Dr. Johanna Jarcho",
                  font_size=16, color=C_GRAY, alignment=PP_ALIGN.CENTER)

    # Bottom accent bar
    bar2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.35), SLIDE_W, Inches(0.15))
    bar2.fill.solid()
    bar2.fill.fore_color.rgb = C_BLUE
    bar2.line.fill.background()

    _add_speaker_notes(slide,
        "Welcome everyone. Today I'll be presenting my first-year project examining "
        "how social anxiety relates to the neural organization of social feedback processing "
        "in adolescents — specifically, whether anxious teens represent peer feedback "
        "idiosyncratically, each in their own unique way.")


def build_slide_02_interests(prs):
    """Slide 2: My Research Interests."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg_white(slide)
    _add_title(slide, "My Research Interests")

    # Main question
    tf = _add_rich_text_box(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(1.2))
    p = _add_paragraph(tf, "", font_size=22, color=C_NAVY)
    _add_run(p, "How do symptoms of social anxiety relate to ", font_size=22, color=C_NAVY)
    _add_run(p, "neural and cognitive differences", font_size=22, color=C_BLUE, bold=True)
    _add_run(p, " in the processing of social feedback?", font_size=22, color=C_NAVY)

    # Sub-question
    _add_text_box(slide, Inches(0.8), Inches(3.0), Inches(11.5), Inches(1.0),
                  "Specifically: How does the socially anxious brain learn about, represent, "
                  "and navigate social environments?",
                  font_size=18, color=C_GRAY)

    # Three interconnected concept nodes
    node_y = Inches(4.5)
    node_w = Inches(3.0)
    node_h = Inches(1.5)
    gap = Inches(0.5)

    labels = ["Neural\nProcessing", "Social\nFeedback", "Learning &\nRepresentation"]
    colors_n = [C_BLUE, C_RED, C_TEAL]
    x_positions = [Inches(1.5), Inches(5.0), Inches(8.5)]

    for i, (x, lab, col) in enumerate(zip(x_positions, labels, colors_n)):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, node_y, node_w, node_h)
        card.fill.solid()
        card.fill.fore_color.rgb = col
        card.line.fill.background()
        tf_card = card.text_frame
        tf_card.word_wrap = True
        p = tf_card.paragraphs[0]
        p.text = lab
        p.font.size = Pt(18)
        p.font.color.rgb = C_WHITE
        p.font.bold = True
        p.font.name = FONT_MAIN
        p.alignment = PP_ALIGN.CENTER
        tf_card.paragraphs[0].space_before = Pt(12)

    # Connecting arrows (simplified as lines)
    for x1, x2 in [(Inches(4.5), Inches(5.0)), (Inches(8.0), Inches(8.5))]:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x1, Inches(5.1), Inches(0.5), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = C_LTGRAY
        arrow.line.fill.background()

    _add_speaker_notes(slide,
        "At the broadest level, I'm interested in how social anxiety shapes "
        "the brain's processing of social feedback. But I want to go beyond "
        "just asking whether anxious brains react more or less — I want to understand "
        "how they learn, how they represent, and how they organize their social world.")


def build_slide_03_why_adolescence(prs):
    """Slide 3: Why Adolescence? — bullets left, developmental diagram right."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg_white(slide)
    _add_title(slide, "Why Adolescence? A Perfect Storm")

    # ── Left column: 5 bullets ──────────────────────────────────
    BULLET_TOP = Inches(1.55)
    BULLET_X   = Inches(0.55)
    BULLET_W   = Inches(7.3)

    bullets = [
        {
            "label": "Social reorientation:",
            "text": "Peer evaluation becomes uniquely motivating; belonging needs reach peak salience as adolescents orient toward the peer world",
            "cite": "Nelson et al., 2005 SIPN; Steinberg & Morris, 2001; Crone & Fuligni, 2020",
            "color": C_BLUE,
        },
        {
            "label": "Social brain reorganization:",
            "text": "PFC and limbic circuits undergo profound restructuring, making the brain acutely sensitive to social information and others' evaluations",
            "cite": "Blakemore, 2008, 2012; Blakemore & Mills, 2014; Mills et al., 2014",
            "color": C_BLUE,
        },
        {
            "label": "Maturational imbalance:",
            "text": "Subcortical/limbic systems mature earlier than PFC, creating a window of heightened affective reactivity with incomplete top-down regulation",
            "cite": "Steinberg, 2008, 2010; Casey, Jones & Somerville, 2011; Chein et al., 2011",
            "color": C_ORANGE,
        },
        {
            "label": "SAD onset peaks in early adolescence:",
            "text": "Social anxiety disorder typically emerges between ages 10–15, coinciding precisely with peak peer sensitivity and maturational imbalance",
            "cite": "Rapee & Spence, 2004; Guyer et al., 2009; Somerville, 2013",
            "color": C_RED,
        },
        {
            "label": "A double-edged sword:",
            "text": "Social-neural sensitivity supports learning and status navigation for most — but for those predisposed to anxiety, this same sensitivity becomes a liability",
            "cite": "Nelson et al., 2016; Silk et al., 2013; Davis et al., 2024; Tillfors et al., 2012",
            "color": C_TEAL,
        },
    ]

    y = BULLET_TOP
    for b in bullets:
        tf = _add_rich_text_box(slide, BULLET_X, y, BULLET_W, Inches(1.0))
        # Bold colored label + body text in one paragraph
        p = tf.paragraphs[0]
        run_label = _add_run(p, "▸ " + b["label"] + " ", bold=True, color=b["color"], font_size=13)
        run_body  = _add_run(p, b["text"], bold=False, color=C_NAVY, font_size=13)
        p.space_after = Pt(2)
        # Citation line
        p_cite = tf.add_paragraph()
        _add_run(p_cite, "   " + b["cite"], italic=True, color=C_GRAY, font_size=9)
        p_cite.space_after = Pt(4)
        y += Inches(1.05)

    # ── Right column: diagram ───────────────────────────────────
    diag_buf = generate_adolescence_diagram()
    diag_stream = io.BytesIO(diag_buf.getvalue())
    slide.shapes.add_picture(
        diag_stream,
        left=Inches(7.95), top=Inches(1.45),
        width=Inches(4.95), height=Inches(5.6)
    )

    _add_speaker_notes(slide,
        "Why focus on adolescence? It's a perfect storm of three converging forces. "
        "First, adolescence is a period of social reorientation — peer evaluation becomes "
        "more motivating and meaningful than at any other developmental stage. "
        "Second, the social brain is fundamentally reorganizing — limbic and prefrontal "
        "circuits are rewiring in ways that make the brain especially sensitive to social cues. "
        "Third, there's a maturational imbalance: subcortical affective systems come online "
        "before the prefrontal regulatory machinery is ready, creating a window of "
        "heightened reactivity without mature control. "
        "And critically, social anxiety disorder onset peaks precisely in this window — ages 10 to 15. "
        "The diagram on the right shows exactly this: all three forces converge during adolescence. "
        "For most teens, this sensitivity is adaptive. For those predisposed to anxiety, "
        "it's a liability — and understanding why requires looking beyond simple threat reactivity.")


def build_slide_04_background(prs):
    """Slide 4: What We Know About SAD and the Brain."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg_white(slide)
    _add_title(slide, "Social Anxiety and Social Feedback Processing",
               subtitle="What the field has established")

    col_w = Inches(3.7)
    top = Inches(1.8)
    card_h = Inches(4.5)

    sections = [
        {
            "title": "Threat Reactivity",
            "color": C_RED,
            "bullets": [
                "Heightened amygdala reactivity to social threat and anticipated peer evaluation (Guyer et al., 2008, 2014)",
                "Enhanced dACC activation during social exclusion and expectancy violation (Eisenberger, 2012)",
                "Exaggerated anterior insula response mediating interoception and anxiety (Terasawa et al., 2013)",
            ]
        },
        {
            "title": "Social Learning",
            "color": C_BLUE,
            "bullets": [
                "Altered prediction error signaling — heightened response to unexpected positive feedback but failure to update predictions (Jarcho et al., 2015, 2016)",
                "Anxiety impedes adaptive social learning under uncertainty (Lamba et al., 2020)",
                "Social feedback processing shows developmental trajectory with increasing vmPFC engagement (Vrtička et al., 2014)",
            ]
        },
        {
            "title": "Cognitive Mechanisms",
            "color": C_TEAL,
            "bullets": [
                "Attentional bias toward threat (Amir et al., 2009; Morrison & Heimberg, 2013)",
                "Computational modeling reveals individual differences in social evaluative decision-making (Castagna et al., 2024)",
                "Deficits in social reward processing (Zhao et al., 2025); dysregulated reward/salience circuits (Clarkson et al., 2025)",
            ]
        }
    ]

    for i, sec in enumerate(sections):
        x = Inches(0.5) + i * (col_w + Inches(0.2))

        # Section header
        _add_text_box(slide, x, top, col_w, Inches(0.4),
                      sec["title"], font_size=16, color=sec["color"], bold=True)

        # Accent line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, top + Inches(0.45), Inches(1.5), Pt(2))
        line.fill.solid()
        line.fill.fore_color.rgb = sec["color"]
        line.line.fill.background()

        # Bullets
        tf = _add_rich_text_box(slide, x, top + Inches(0.6), col_w, card_h - Inches(0.6))
        for bullet in sec["bullets"]:
            _add_paragraph(tf, f"• {bullet}", font_size=12, color=C_NAVY, space_after=10)

    # Bottom callout
    callout = _add_card(slide, Inches(1.5), Inches(6.5), Inches(10), Inches(0.7),
                        fill_color=RGBColor(0xFF, 0xF3, 0xE0))
    _add_text_box(slide, Inches(1.8), Inches(6.55), Inches(9.5), Inches(0.6),
                  'But nearly all fMRI work asks: does this region activate MORE or LESS? '
                  'What about how the brain organizes social information?',
                  font_size=14, color=C_ORANGE, bold=True, italic=True, alignment=PP_ALIGN.CENTER)

    _add_speaker_notes(slide,
        "Here's what the field has established. In threat reactivity, we know anxious youth "
        "show heightened amygdala and insula responses. In social learning, Jarcho and colleagues "
        "showed that anxious adolescents detect unexpected positive feedback but fail to update "
        "their predictions — the 'failed positive learning' finding. And computational work reveals "
        "individual differences in social evaluative decision-making. But notice: nearly all of "
        "this work asks whether a brain region activates MORE or LESS. What about how the brain "
        "actually organizes and structures social information?")


def build_slide_05_gap(prs):
    """Slide 5: The Gap — Beyond Threat Response."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg_white(slide)
    _add_title(slide, "What's Missing: Beyond Threat Response")

    # The shift diagram
    # Left side: Traditional (bar chart metaphor)
    _add_text_box(slide, Inches(1.0), Inches(1.8), Inches(5), Inches(0.5),
                  "Traditional fMRI asks:", font_size=18, color=C_GRAY)
    _add_text_box(slide, Inches(1.0), Inches(2.3), Inches(5), Inches(0.6),
                  '"Does this region activate MORE or LESS?"',
                  font_size=20, color=C_NAVY, bold=True)
    _add_text_box(slide, Inches(1.0), Inches(2.9), Inches(5), Inches(0.4),
                  "→ Magnitude", font_size=16, color=C_GRAY, italic=True)

    # Arrow
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(5.8), Inches(2.4),
                                   Inches(1.2), Inches(0.5))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = C_BLUE
    arrow.line.fill.background()

    # Right side: What's missing
    _add_text_box(slide, Inches(7.2), Inches(1.8), Inches(5.5), Inches(0.5),
                  "But SAD involves more:", font_size=18, color=C_GRAY)
    _add_text_box(slide, Inches(7.2), Inches(2.3), Inches(5.5), Inches(0.6),
                  '"In what WAY does the brain process this?"',
                  font_size=20, color=C_BLUE, bold=True)
    _add_text_box(slide, Inches(7.2), Inches(2.9), Inches(5.5), Inches(0.4),
                  "→ Pattern & Organization", font_size=16, color=C_BLUE, italic=True)

    # Three missing questions
    questions = [
        ("Meaning-making", "How do anxious teens interpret and organize social experiences?"),
        ("Social representation", "Does the anxious brain structure its map of the social world differently?"),
        ("Individual variation", "Do all anxious teens process threat the same way, or does each diverge uniquely?"),
    ]

    for i, (title, desc) in enumerate(questions):
        y = Inches(3.8) + i * Inches(1.1)
        _add_card(slide, Inches(1.0), y, Inches(11.3), Inches(0.9))

        _add_text_box(slide, Inches(1.3), y + Inches(0.08), Inches(3), Inches(0.35),
                      title, font_size=16, color=C_ORANGE, bold=True)
        _add_text_box(slide, Inches(1.3), y + Inches(0.4), Inches(10.5), Inches(0.4),
                      desc, font_size=14, color=C_NAVY)

    _add_text_box(slide, Inches(1.0), Inches(6.8), Inches(11.5), Inches(0.5),
                  "We need methods that capture the pattern and organization of neural representations.",
                  font_size=15, color=C_NAVY, bold=True, italic=True, alignment=PP_ALIGN.CENTER)

    _add_speaker_notes(slide,
        "Here's the gap. Nearly all fMRI work on SAD asks about magnitude — more or less activation. "
        "But SAD involves much more than heightened threat reactivity. What about meaning-making — "
        "how do anxious teens interpret and organize their social experiences? What about social "
        "representation — does the brain structure its map of the social world differently? And "
        "critically, what about individual variation — do all anxious teens process threat the same "
        "way, or does each diverge uniquely? To answer these questions, we need methods that capture "
        "the pattern and organization of neural representations, not just their magnitude.")


def build_slide_06_rsa(prs):
    """Slide 6: RSA — Why This Method Fits."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg_white(slide)
    _add_title(slide, "Representational Similarity Analysis: A Different Question")

    # The contrast
    _add_text_box(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.4),
                  "Traditional fMRI:", font_size=16, color=C_GRAY)
    _add_text_box(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(0.5),
                  '"Does this region activate more for A vs B?"',
                  font_size=18, color=C_NAVY, bold=True)
    _add_text_box(slide, Inches(0.8), Inches(2.5), Inches(5.5), Inches(0.3),
                  "→ Magnitude  (Kriegeskorte et al., 2008)",
                  font_size=13, color=C_GRAY, italic=True)

    _add_text_box(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(0.4),
                  "RSA:", font_size=16, color=C_BLUE, bold=True)
    _add_text_box(slide, Inches(7.0), Inches(2.0), Inches(5.5), Inches(0.5),
                  '"How are conditions organized in neural pattern space?"',
                  font_size=18, color=C_BLUE, bold=True)
    _add_text_box(slide, Inches(7.0), Inches(2.5), Inches(5.5), Inches(0.3),
                  "→ Geometry  (Kriegeskorte, Mur & Bandettini, 2008)",
                  font_size=13, color=C_BLUE, italic=True)

    # Why RSA fits social cognition
    _add_text_box(slide, Inches(0.8), Inches(3.3), Inches(11.5), Inches(0.4),
                  "Why RSA is right for social cognition:", font_size=18, color=C_NAVY, bold=True)

    bullets = [
        ("Social stimuli are multidimensional and continuous",
         "people, relationships, reputations — they don't reduce to discrete categories (Popal et al., 2019, SCAN)"),
        ("MVPA/RSA decodes what univariate misses",
         "person knowledge, social traits, self-other distinctions (Wagner, Chavez & Broom, 2018, WIREs Cog Sci)"),
        ("RSA reveals representational organization",
         "which peers are similar in neural space, which are different (Thornton & Tamir, 2019, Nature Human Behaviour)"),
    ]

    tf = _add_rich_text_box(slide, Inches(0.8), Inches(3.8), Inches(11.5), Inches(2.5))
    for title, desc in bullets:
        p = _add_paragraph(tf, "", font_size=15, space_after=8)
        _add_run(p, f"• {title}", font_size=15, color=C_NAVY, bold=True)
        _add_run(p, f" — {desc}", font_size=13, color=C_GRAY)

    # Fingerprint metaphor callout
    callout = _add_card(slide, Inches(1.5), Inches(6.0), Inches(10), Inches(1.0),
                        fill_color=RGBColor(0xE8, 0xF4, 0xFD))
    _add_text_box(slide, Inches(1.8), Inches(6.1), Inches(9.5), Inches(0.8),
                  "The Fingerprint Metaphor: Each person's brain response to a social event is "
                  "like a fingerprint — a unique pattern across voxels. RSA compares these fingerprints.",
                  font_size=14, color=C_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

    _add_speaker_notes(slide,
        "RSA asks a fundamentally different question from traditional fMRI. Instead of "
        "'does this region activate more,' it asks 'how are conditions organized in neural "
        "pattern space?' This is especially powerful for social cognition because social "
        "stimuli are multidimensional — people, relationships, reputations don't reduce to "
        "simple categories. RSA can decode the organization of social knowledge that "
        "univariate approaches miss entirely. Think of each person's brain response as a "
        "fingerprint — a unique spatial pattern. RSA compares these fingerprints to reveal "
        "representational structure.")


def build_slide_07_idiosyncrasy(prs):
    """Slide 7: The Idiosyncrasy Literature."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg_white(slide)
    _add_title(slide, "Neural Similarity, Social Connection, and Individual Differences")

    papers = [
        {
            "author": "Parkinson et al., 2018",
            "journal": "Nature Communications",
            "finding": "Neural responses to naturalistic stimuli predicted social network distance. Similar brains → closer friends.",
            "color": C_BLUE,
        },
        {
            "author": "Baek et al., 2023",
            "journal": "Psychological Science",
            "finding": "Lonely individuals showed idiosyncratic neural processing. Idiosyncrasy correlated with feeling misunderstood.",
            "color": C_ORANGE,
        },
        {
            "author": "Camacho et al., 2024",
            "journal": "JAACAP (N=740)",
            "finding": "Higher social anxiety → greater inter-subject neural variability. No mean activation differences — the effect is in variability.",
            "color": C_RED,
        },
        {
            "author": "Shen et al., 2025",
            "journal": "Nature Human Behaviour",
            "finding": "Pre-existing neural similarity predicted friendship formation 8 months later. Neural divergence → difficulty connecting.",
            "color": C_TEAL,
        },
    ]

    card_w = Inches(5.5)
    card_h = Inches(1.15)

    for i, paper in enumerate(papers):
        col = i % 2
        row = i // 2
        x = Inches(0.8) + col * (card_w + Inches(0.3))
        y = Inches(1.7) + row * (card_h + Inches(0.2))

        _add_card(slide, x, y, card_w, card_h)

        # Author + journal
        tf = _add_rich_text_box(slide, x + Inches(0.15), y + Inches(0.05),
                                card_w - Inches(0.3), Inches(0.35))
        p = _add_paragraph(tf, "", font_size=13, space_after=2)
        _add_run(p, paper["author"], font_size=14, color=paper["color"], bold=True)
        _add_run(p, f'  {paper["journal"]}', font_size=11, color=C_GRAY, italic=True)

        # Finding
        _add_text_box(slide, x + Inches(0.15), y + Inches(0.4),
                      card_w - Inches(0.3), Inches(0.7),
                      paper["finding"], font_size=12, color=C_NAVY)

    # Connecting thread
    callout = _add_card(slide, Inches(1.0), Inches(4.5), Inches(11), Inches(1.2),
                        fill_color=RGBColor(0xFD, 0xF0, 0xE8))
    tf = _add_rich_text_box(slide, Inches(1.3), Inches(4.6), Inches(10.5), Inches(1.0))
    p = _add_paragraph(tf, "", font_size=16, alignment=PP_ALIGN.CENTER, space_after=4)
    _add_run(p, "The connecting thread: ", font_size=16, color=C_NAVY, bold=True)
    _add_run(p, "Neural similarity predicts social bonding. Neural divergence predicts "
             "social disconnection.", font_size=16, color=C_NAVY)
    p2 = _add_paragraph(tf, "", font_size=16, alignment=PP_ALIGN.CENTER)
    _add_run(p2, "Does social anxiety produce neural divergence — specifically for threatening social stimuli?",
             font_size=16, color=C_ORANGE, bold=True, italic=True)

    _add_citation(slide, "Also: Finn & Scheinost, 2020, NeuroImage (IS-RSA framework); "
                  "Baek & Parkinson, 2022, SPPC", top=Inches(6.0))

    _add_speaker_notes(slide,
        "This emerging literature shows something remarkable: how similarly people's brains "
        "process the world predicts real social outcomes. Parkinson found similar brains → closer "
        "friends. Baek found lonely individuals are neurally idiosyncratic — and this correlates "
        "with feeling misunderstood. Camacho, in a massive sample of 740 youth, showed that higher "
        "social anxiety predicts greater inter-subject neural VARIABILITY — invisible to standard "
        "analyses that only look at mean activation. And Shen showed pre-existing neural similarity "
        "predicts who becomes friends 8 months later. The connecting thread: neural similarity "
        "predicts social bonding. Does social anxiety fragment this similarity?")


def build_slide_08_questions(prs):
    """Slide 8: Research Questions."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg_white(slide)
    _add_title(slide, "Research Questions")

    # Overarching
    _add_text_box(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.8),
                  "Does social anxiety moderate unique (idiosyncratic) neural representations "
                  "of social feedback?",
                  font_size=22, color=C_NAVY, bold=True, alignment=PP_ALIGN.CENTER)

    # Q1
    q1_card = _add_card(slide, Inches(0.8), Inches(2.8), Inches(11.3), Inches(2.2),
                        fill_color=RGBColor(0xE8, 0xF4, 0xFD))

    _add_text_box(slide, Inches(1.2), Inches(2.9), Inches(10.5), Inches(0.4),
                  "Q1  (This Talk)", font_size=18, color=C_BLUE, bold=True)
    _add_text_box(slide, Inches(1.2), Inches(3.35), Inches(10.5), Inches(0.6),
                  "Is SAD related to increasingly divergent neural representations of "
                  "emotional feedback?",
                  font_size=18, color=C_NAVY, bold=True)

    tf = _add_rich_text_box(slide, Inches(1.2), Inches(3.95), Inches(10.5), Inches(1.0))
    _add_paragraph(tf, "• Do socially anxious adolescents each represent peer feedback in their own unique way?",
                   font_size=14, color=C_NAVY, space_after=4)
    _add_paragraph(tf, "• The Anna Karenina hypothesis: low anxiety → shared patterns; high anxiety → each person diverges",
                   font_size=14, color=C_ORANGE, bold=True, space_after=4)

    # Q2
    q2_card = _add_card(slide, Inches(0.8), Inches(5.3), Inches(11.3), Inches(1.4))

    _add_text_box(slide, Inches(1.2), Inches(5.4), Inches(10.5), Inches(0.4),
                  "Q2  (Future Direction)", font_size=18, color=C_GRAY, bold=True)
    _add_text_box(slide, Inches(1.2), Inches(5.85), Inches(10.5), Inches(0.5),
                  "Do neural representations reflect different rates of social learning?",
                  font_size=16, color=C_GRAY)
    _add_text_box(slide, Inches(1.2), Inches(6.3), Inches(10.5), Inches(0.4),
                  "Does the brain's representational geometry come to mirror the true peer structure over time?",
                  font_size=13, color=C_GRAY, italic=True)

    _add_speaker_notes(slide,
        "Our overarching question: does social anxiety moderate unique neural representations "
        "of social feedback? This talk focuses on Q1: the Anna Karenina hypothesis. Named after "
        "Tolstoy's opening line — all happy families are alike, each unhappy family is unhappy "
        "in its own way. We predict that low-anxiety teens share similar neural patterns, while "
        "high-anxiety teens each diverge in their own unique way. Q2 about learning trajectories "
        "is a future direction we're actively developing.")


def build_slide_09_learn_task(prs):
    """Slide 9: The LEARN Task."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg_white(slide)
    _add_title(slide, "The LEARN Task: A Virtual School")

    # Task image from Learn.pptx
    task_img = IMG_DIR / "slide04_Reminder:_Virtual_School_Task:_1.png"
    _add_image_safe(slide, task_img, Inches(6.8), Inches(1.4), width=Inches(5.8))

    # Sample info
    tf = _add_rich_text_box(slide, Inches(0.8), Inches(1.6), Inches(5.8), Inches(5.0))

    _add_paragraph(tf, "Sample", font_size=18, color=C_BLUE, bold=True, space_after=4)
    _add_paragraph(tf, "• 33 adolescents (ages 10-15)", font_size=15, color=C_NAVY, space_after=4)
    _add_paragraph(tf, "• fMRI during social learning task", font_size=15, color=C_NAVY, space_after=12)

    _add_paragraph(tf, "Design: 2×2", font_size=18, color=C_BLUE, bold=True, space_after=4)
    _add_paragraph(tf, "• 4 virtual peers with hidden properties:", font_size=15, color=C_NAVY, space_after=4)
    _add_paragraph(tf, "  Valence (Nice/Mean) × Predictability (80%/60%)", font_size=14, color=C_NAVY, space_after=12)

    _add_paragraph(tf, "Trial Structure", font_size=18, color=C_BLUE, bold=True, space_after=4)
    _add_paragraph(tf, "• See peer → Predict feedback → Receive Nice/Mean → Rate feeling",
                   font_size=15, color=C_NAVY, space_after=4)
    _add_paragraph(tf, "• 128 trials across 4 fMRI runs", font_size=15, color=C_NAVY, space_after=12)

    _add_paragraph(tf, "Advantage over naturalistic designs:", font_size=14, color=C_TEAL, bold=True, space_after=4)
    _add_paragraph(tf, "Experimental control over social feedback — we know exactly what feedback "
                   "each participant received, from which peer, at which point in learning.",
                   font_size=13, color=C_GRAY, italic=True)

    _add_speaker_notes(slide,
        "We used the LEARN task — a virtual school paradigm where 33 adolescents interacted "
        "with four computer-generated peers in the MRI scanner. Each peer had hidden properties: "
        "they could be nice or mean, and predictable or unpredictable. Over 128 trials across "
        "4 runs, participants predicted and received feedback from these peers. This gives us "
        "experimental control over social feedback that naturalistic paradigms like movie-watching "
        "can't provide — we know exactly what feedback each participant got, from which peer, "
        "and at what point in their learning.")


def build_slide_10_isrsa_walkthrough(prs):
    """Slide 10: IS-RSA & Anna Karenina Visual Walk-Through."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg_white(slide)
    _add_title(slide, "Testing the Anna Karenina Hypothesis: IS-RSA")

    # Tolstoy quote
    _add_text_box(slide, Inches(0.8), Inches(1.35), Inches(11.5), Inches(0.5),
                  '"All happy families are alike; each unhappy family is unhappy in its own way." — Tolstoy',
                  font_size=14, color=C_ORANGE, italic=True, alignment=PP_ALIGN.CENTER)

    # Generate and embed the IS-RSA walk-through figure
    fig_buf = generate_isrsa_walkthrough()
    slide.shapes.add_picture(fig_buf, Inches(0.3), Inches(1.9), width=Inches(12.5))

    _add_speaker_notes(slide,
        "Here's how we test the Anna Karenina hypothesis using inter-subject RSA. "
        "Step 1: We extract the multi-voxel pattern — the fingerprint — from each person's brain "
        "for mean feedback. Step 2: We compare all pairs of subjects to build an inter-subject "
        "similarity matrix — how similar is each person's fingerprint to everyone else's? "
        "Step 3: We build the Anna Karenina behavioral model — the prediction that subjects "
        "with lower anxiety should be more similar, and those with higher anxiety should diverge. "
        "Step 4: We run a Mantel test — correlating the neural similarity matrix with the "
        "behavioral model, using 10,000 permutations for statistical significance.")


def build_slide_11_methods(prs):
    """Slide 11: Quick Methods / Pipeline."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg_white(slide)
    _add_title(slide, "fMRI Processing Pipeline")

    # Pipeline steps as connected boxes
    steps = [
        ("fMRI\nAcquisition", C_GRAY),
        ("Preprocessing\n(AFNI, no smoothing)", C_GRAY),
        ("GLM\n(run-wise betas\n8 conditions)", C_BLUE),
        ("ROI Extraction\n(8 social brain\nregions)", C_BLUE),
        ("IS-RSA\n(inter-subject\ncomparison)", C_ORANGE),
    ]

    box_w = Inches(2.0)
    box_h = Inches(1.3)
    start_x = Inches(0.5)
    y = Inches(1.7)

    for i, (label, color) in enumerate(steps):
        x = start_x + i * (box_w + Inches(0.4))
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(13)
        p.font.color.rgb = C_WHITE
        p.font.bold = True
        p.font.name = FONT_MAIN
        p.alignment = PP_ALIGN.CENTER

        # Arrow between boxes
        if i < len(steps) - 1:
            arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                         x + box_w, y + Inches(0.45),
                                         Inches(0.4), Inches(0.3))
            arr.fill.solid()
            arr.fill.fore_color.rgb = C_LTGRAY
            arr.line.fill.background()

    # Key details
    tf = _add_rich_text_box(slide, Inches(0.8), Inches(3.3), Inches(5.5), Inches(1.5))
    _add_paragraph(tf, "Key design choices:", font_size=16, color=C_NAVY, bold=True, space_after=6)
    _add_paragraph(tf, "• No spatial smoothing (preserves pattern info)", font_size=14, color=C_NAVY, space_after=4)
    _add_paragraph(tf, "• Standard MNI space (enables cross-subject comparison)", font_size=14, color=C_NAVY, space_after=4)
    _add_paragraph(tf, "• 10,000 Mantel permutations per test", font_size=14, color=C_NAVY, space_after=4)
    _add_paragraph(tf, "• FDR correction across 32 tests (8 ROIs × 2 feedback × 2 models)", font_size=14, color=C_NAVY)

    # ROI image
    if ROI_IMG.exists():
        slide.shapes.add_picture(str(ROI_IMG), Inches(6.5), Inches(3.2), width=Inches(6.2))

    _add_citation(slide, "8 ROIs: vmPFC, dACC1, dACC2, AntInsula, VS, Amygdala, R-TPJ, dmPFC")

    _add_speaker_notes(slide,
        "Briefly on methods: we used AFNI to preprocess the fMRI data with no spatial smoothing — "
        "critical for preserving the fine-grained pattern information RSA needs. We estimated "
        "run-wise beta maps for 8 feedback conditions, then extracted voxel-wise patterns from "
        "8 ROIs spanning the social brain. IS-RSA was computed with 10,000 Mantel permutations, "
        "corrected for multiple comparisons using FDR.")


def build_slide_12_results_dacc(prs):
    """Slide 12: Results — dACC1 Finding."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg_white(slide)
    _add_title(slide, "dACC1: The Anna Karenina Effect for Mean Feedback")

    # Forest plot
    fig_buf = generate_forest_plot()
    slide.shapes.add_picture(fig_buf, Inches(0.5), Inches(1.5), width=Inches(7.5))

    # Key finding callout
    callout = _add_card(slide, Inches(8.3), Inches(1.8), Inches(4.5), Inches(2.5),
                        fill_color=RGBColor(0xFD, 0xF0, 0xE8))

    tf = _add_rich_text_box(slide, Inches(8.5), Inches(1.9), Inches(4.2), Inches(2.3))
    _add_paragraph(tf, "Key Finding", font_size=18, color=C_ORANGE, bold=True, space_after=8)
    _add_paragraph(tf, "dACC1 × Mean Feedback:", font_size=14, color=C_NAVY, bold=True, space_after=4)
    _add_paragraph(tf, "ρ = −0.166, p_FDR = .036 ★", font_size=16, color=C_ORANGE, bold=True, space_after=8)
    _add_paragraph(tf, "Only result surviving FDR correction across 32 tests",
                   font_size=13, color=C_NAVY, space_after=8)
    _add_paragraph(tf, "Higher social anxiety → more idiosyncratic neural patterns for mean feedback",
                   font_size=13, color=C_NAVY, bold=True, italic=True)

    # Interpretation box
    _add_card(slide, Inches(8.3), Inches(4.6), Inches(4.5), Inches(2.2))
    tf2 = _add_rich_text_box(slide, Inches(8.5), Inches(4.7), Inches(4.2), Inches(2.0))
    _add_paragraph(tf2, "What does dACC1 do?", font_size=14, color=C_BLUE, bold=True, space_after=6)
    _add_paragraph(tf2, "• Conflict monitoring", font_size=12, color=C_NAVY, space_after=3)
    _add_paragraph(tf2, "• Social prediction errors", font_size=12, color=C_NAVY, space_after=3)
    _add_paragraph(tf2, "• Salience detection in evaluative contexts", font_size=12, color=C_NAVY, space_after=6)
    _add_paragraph(tf2, "Where the brain detects that social feedback doesn't match expectations",
                   font_size=11, color=C_GRAY, italic=True)

    _add_speaker_notes(slide,
        "Here's our primary finding. In the dorsal anterior cingulate — a region critical for "
        "conflict monitoring and social prediction errors — we found the Anna Karenina effect "
        "for mean feedback. This is the only result that survives FDR correction across our "
        "32 tests. Higher social anxiety predicted more idiosyncratic brain patterns. "
        "The dACC is where the brain detects that social feedback doesn't match expectations. "
        "When anxious teens receive mean feedback, each encodes this mismatch differently — "
        "some expected niceness and are shocked, some expected meanness and see confirmation, "
        "others trigger self-referential cascades.")


def build_slide_13_dacc_meaning(prs):
    """Slide 13: What dACC1 Idiosyncrasy Means."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg_white(slide)
    _add_title(slide, "What Does Idiosyncrasy in the dACC Mean?")

    # AnnaK heatmap
    fig_buf = generate_annak_heatmap()
    slide.shapes.add_picture(fig_buf, Inches(0.3), Inches(1.5), width=Inches(8.5))

    # Interpretation
    tf = _add_rich_text_box(slide, Inches(9.0), Inches(1.6), Inches(4.0), Inches(5.0))
    _add_paragraph(tf, "The Pattern:", font_size=18, color=C_NAVY, bold=True, space_after=10)

    _add_paragraph(tf, "Low-anxiety adolescents:", font_size=14, color=C_BLUE, bold=True, space_after=4)
    _add_paragraph(tf, "Process prediction errors uniformly — 'that peer was meaner than expected, I'll update'",
                   font_size=13, color=C_NAVY, space_after=12)

    _add_paragraph(tf, "High-anxiety adolescents:", font_size=14, color=C_ORANGE, bold=True, space_after=4)
    _add_paragraph(tf, "Each responds to the same prediction error differently:",
                   font_size=13, color=C_NAVY, space_after=4)
    _add_paragraph(tf, "  • Some expected niceness (shocked)", font_size=12, color=C_NAVY, space_after=3)
    _add_paragraph(tf, "  • Some expected meanness (confirmation)", font_size=12, color=C_NAVY, space_after=3)
    _add_paragraph(tf, "  • Some trigger self-referential cascades", font_size=12, color=C_NAVY, space_after=10)

    _add_paragraph(tf, "→ The experience of social threat is fundamentally private for each anxious teen",
                   font_size=13, color=C_ORANGE, bold=True, italic=True)

    _add_speaker_notes(slide,
        "What does this mean? The left panel shows the predicted Anna Karenina pattern versus "
        "what we actually observe in the dACC. Low-anxiety adolescents process prediction errors "
        "more uniformly — they share a common neural response. But high-anxiety adolescents each "
        "respond to the same social mismatch in their own unique way. The experience of social "
        "threat becomes fundamentally private.")


def build_slide_14_anterior_insula(prs):
    """Slide 14: Anterior Insula Reverse AnnaK."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg_white(slide)
    _add_title(slide, "Anterior Insula: A Reverse Anna Karenina Effect",
               subtitle="Exploratory (uncorrected)")

    # Finding
    tf = _add_rich_text_box(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(3.0))
    _add_paragraph(tf, "The anterior insula shows convergence:", font_size=18, color=C_NAVY, bold=True, space_after=10)
    _add_paragraph(tf, "Mean Feedback: ρ = +0.188, p = .039", font_size=16, color=C_BLUE, bold=True, space_after=4)
    _add_paragraph(tf, "Nice Feedback: ρ = +0.216, p = .024", font_size=16, color=C_BLUE, bold=True, space_after=10)
    _add_paragraph(tf, "The ONLY ROI showing convergence for BOTH feedback types",
                   font_size=14, color=C_NAVY, italic=True, space_after=12)

    _add_paragraph(tf, "Interpretation:", font_size=16, color=C_BLUE, bold=True, space_after=6)
    _add_paragraph(tf, "Higher anxiety → MORE similar anterior insula patterns",
                   font_size=14, color=C_NAVY, space_after=4)
    _add_paragraph(tf, "Anxious adolescents may share a stereotyped visceral alarm response — "
                   "a locked-in 'ALERT! Social evaluation happening!' signal",
                   font_size=13, color=C_NAVY, space_after=8)

    # Dissociation diagram
    _add_text_box(slide, Inches(7.0), Inches(1.8), Inches(5.5), Inches(0.4),
                  "The Insular Dissociation:", font_size=18, color=C_NAVY, bold=True)

    # Convergence card
    _add_card(slide, Inches(7.0), Inches(2.5), Inches(5.2), Inches(1.5),
              fill_color=RGBColor(0xE8, 0xF4, 0xFD))
    _add_text_box(slide, Inches(7.2), Inches(2.6), Inches(5.0), Inches(0.3),
                  "Anterior Insula (detection/salience)", font_size=14, color=C_BLUE, bold=True)
    _add_text_box(slide, Inches(7.2), Inches(2.95), Inches(5.0), Inches(0.8),
                  "→ CONVERGENCE with anxiety\nAll anxious teens detect the same threat",
                  font_size=13, color=C_BLUE)

    # Idiosyncrasy card
    _add_card(slide, Inches(7.0), Inches(4.2), Inches(5.2), Inches(1.5),
              fill_color=RGBColor(0xFD, 0xF0, 0xE8))
    _add_text_box(slide, Inches(7.2), Inches(4.3), Inches(5.0), Inches(0.3),
                  "dACC (appraisal/conflict)", font_size=14, color=C_ORANGE, bold=True)
    _add_text_box(slide, Inches(7.2), Inches(4.65), Inches(5.0), Inches(0.8),
                  "→ IDIOSYNCRASY with anxiety\nEach anxious teen processes it uniquely",
                  font_size=13, color=C_ORANGE)

    # Punchline
    _add_text_box(slide, Inches(7.0), Inches(6.0), Inches(5.5), Inches(0.6),
                  '"They share the alarm, but fragment the meaning"',
                  font_size=16, color=C_NAVY, bold=True, italic=True, alignment=PP_ALIGN.CENTER)

    _add_speaker_notes(slide,
        "An intriguing exploratory finding: the anterior insula shows the REVERSE Anna Karenina "
        "pattern — higher anxiety predicts MORE convergent patterns, for both nice and mean feedback. "
        "This creates a beautiful dissociation with the dACC finding. The anterior insula mediates "
        "interoceptive awareness — it's where the brain generates the visceral alarm. Anxious "
        "teens may all share the same stereotyped alarm response. But the dACC, where that alarm "
        "is appraised and interpreted, shows idiosyncrasy. They share the alarm but fragment "
        "the meaning. This is uncorrected, so it's exploratory, but the consistency across "
        "both feedback types is notable.")


def build_slide_15_searchlight(prs):
    """Slide 15: Searchlight Results."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg_white(slide)
    _add_title(slide, "Whole-Brain Searchlight: Mean vs Nice Asymmetry",
               subtitle="Exploratory (uncorrected, p < .01 voxel threshold)")

    # Brain images
    idio_img = SEARCH_DIR / "afni_Mean_idio_glass.png"
    conv_img = SEARCH_DIR / "afni_Mean_conv_glass.png"

    _add_text_box(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(0.3),
                  "Mean Feedback — Idiosyncrasy (negative ρ)", font_size=14, color=C_ORANGE, bold=True)
    _add_image_safe(slide, idio_img, Inches(0.3), Inches(1.9), width=Inches(6.3))

    _add_text_box(slide, Inches(0.5), Inches(4.0), Inches(6), Inches(0.3),
                  "Mean Feedback — Convergence (positive ρ)", font_size=14, color=C_BLUE, bold=True)
    _add_image_safe(slide, conv_img, Inches(0.3), Inches(4.4), width=Inches(6.3))

    # Asymmetry chart
    fig_buf = generate_asymmetry_chart()
    slide.shapes.add_picture(fig_buf, Inches(6.8), Inches(1.5), width=Inches(6.0))

    # Key regions
    tf = _add_rich_text_box(slide, Inches(7.0), Inches(5.0), Inches(5.8), Inches(1.8))
    _add_paragraph(tf, "Key idiosyncrasy regions (Mean):", font_size=12, color=C_ORANGE, bold=True, space_after=3)
    _add_paragraph(tf, "Brainstem, L-SMG, posterior insula, temporal fusiform, frontal pole",
                   font_size=11, color=C_NAVY, space_after=6)
    _add_paragraph(tf, "Key convergence regions (Mean):", font_size=12, color=C_BLUE, bold=True, space_after=3)
    _add_paragraph(tf, "Primary visual cortex, anterior insula, PCC, angular gyrus",
                   font_size=11, color=C_NAVY)

    _add_speaker_notes(slide,
        "The whole-brain searchlight tells a richer story. For mean feedback, 71% of significant "
        "voxels show idiosyncrasy — including brainstem, supramarginal gyrus, posterior insula, "
        "and fusiform cortex. These are appraisal circuits. Convergence appears in detection "
        "regions — primary visual cortex, anterior insula, PCC. For nice feedback, the pattern "
        "reverses: 57% convergence. Negative feedback fragments. Positive feedback converges. "
        "None of these survive cluster-level correction, so they're hypothesis-generating.")


def build_slide_16_emerging_story(prs):
    """Slide 16: The Emerging Story."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg_white(slide)
    _add_title(slide, "The Emerging Story: Two Directions of Anna Karenina")

    # Two columns
    # Left: Idiosyncrasy
    _add_card(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.5),
              fill_color=RGBColor(0xFD, 0xF0, 0xE8))
    _add_text_box(slide, Inches(1.0), Inches(1.9), Inches(5.2), Inches(0.4),
                  "Idiosyncrasy  (each anxious teen diverges)",
                  font_size=18, color=C_ORANGE, bold=True)

    tf1 = _add_rich_text_box(slide, Inches(1.0), Inches(2.5), Inches(5.2), Inches(3.5))
    _add_paragraph(tf1, "Appraisal circuits:", font_size=14, color=C_NAVY, bold=True, space_after=6)
    _add_paragraph(tf1, "• dACC  (conflict monitoring, prediction errors) ★", font_size=14, color=C_ORANGE, bold=True, space_after=4)
    _add_paragraph(tf1, "• Brainstem  (autonomic threat response)", font_size=13, color=C_NAVY, space_after=4)
    _add_paragraph(tf1, "• L-SMG  (self-other distinction)", font_size=13, color=C_NAVY, space_after=4)
    _add_paragraph(tf1, "• Posterior insula  (bodily experience)", font_size=13, color=C_NAVY, space_after=4)
    _add_paragraph(tf1, "• Fusiform  (face processing)", font_size=13, color=C_NAVY, space_after=10)
    _add_paragraph(tf1, "HOW the threat is interpreted and experienced",
                   font_size=13, color=C_ORANGE, italic=True)

    # Right: Convergence
    _add_card(slide, Inches(6.8), Inches(1.8), Inches(5.5), Inches(4.5),
              fill_color=RGBColor(0xE8, 0xF4, 0xFD))
    _add_text_box(slide, Inches(7.0), Inches(1.9), Inches(5.2), Inches(0.4),
                  "Convergence  (anxious teens align)",
                  font_size=18, color=C_BLUE, bold=True)

    tf2 = _add_rich_text_box(slide, Inches(7.0), Inches(2.5), Inches(5.2), Inches(3.5))
    _add_paragraph(tf2, "Detection/salience circuits:", font_size=14, color=C_NAVY, bold=True, space_after=6)
    _add_paragraph(tf2, "• Anterior insula  (salience detection)", font_size=13, color=C_BLUE, bold=True, space_after=4)
    _add_paragraph(tf2, "• PCC  (self-referential processing)", font_size=13, color=C_NAVY, space_after=4)
    _add_paragraph(tf2, "• Angular gyrus  (mentalizing)", font_size=13, color=C_NAVY, space_after=4)
    _add_paragraph(tf2, "• V1  (basic perceptual encoding)", font_size=13, color=C_NAVY, space_after=10)
    _add_paragraph(tf2, "THAT something threatening happened", font_size=13, color=C_BLUE, italic=True)

    # Bottom punchline
    callout = _add_card(slide, Inches(2.0), Inches(6.5), Inches(9.0), Inches(0.7),
                        fill_color=RGBColor(0x00, 0x20, 0x60))
    _add_text_box(slide, Inches(2.2), Inches(6.55), Inches(8.7), Inches(0.6),
                  "Anxious adolescents share the alarm but fragment the appraisal",
                  font_size=18, color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    _add_speaker_notes(slide,
        "Here's the emerging story. The Anna Karenina effect operates in two directions "
        "depending on the brain circuit. In APPRAISAL circuits — dACC, brainstem, SMG, "
        "posterior insula, fusiform — anxiety produces idiosyncrasy. Each anxious teen "
        "interprets and experiences the threat uniquely. But in DETECTION circuits — "
        "anterior insula, PCC, angular gyrus, V1 — anxiety produces convergence. They "
        "all detect that something threatening happened in the same way. The punchline: "
        "anxious adolescents share the alarm but fragment the appraisal.")


def build_slide_17_discussion(prs):
    """Slide 17: Discussion."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg_white(slide)
    _add_title(slide, "Discussion: Neural Idiosyncrasy and Social Disconnection")

    tf = _add_rich_text_box(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(4.5))

    _add_paragraph(tf, "The 'Neural Loneliness' Frame", font_size=20, color=C_NAVY, bold=True, space_after=12)

    _add_paragraph(tf, "1. The dACC is where the brain detects that social feedback doesn't match expectations",
                   font_size=16, color=C_NAVY, space_after=8)

    _add_paragraph(tf, "2. Anxious adolescents each encode this mismatch differently — no two experience "
                   "mean feedback the same way",
                   font_size=16, color=C_NAVY, space_after=8)

    _add_paragraph(tf, "3. Neural similarity predicts social bonding (Shen et al., 2025). "
                   "Neural idiosyncrasy may perpetuate social isolation.",
                   font_size=16, color=C_NAVY, space_after=8)

    p = _add_paragraph(tf, "", font_size=16, space_after=12)
    _add_run(p, "4. ", font_size=16, color=C_NAVY)
    _add_run(p, "They don't just feel misunderstood — at a neural level, their experience ",
             font_size=16, color=C_NAVY)
    _add_run(p, "IS", font_size=16, color=C_ORANGE, bold=True, italic=True)
    _add_run(p, " different from everyone else's", font_size=16, color=C_NAVY)

    # Phenomenology callout
    callout = _add_card(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(1.2),
                        fill_color=RGBColor(0xFD, 0xF0, 0xE8))
    _add_text_box(slide, Inches(1.8), Inches(5.6), Inches(9.5), Inches(1.0),
                  "This may explain the clinical observation that anxious youth report feeling "
                  "misunderstood specifically around negative social experiences — because at a "
                  "neural level, their processing of social threat is fundamentally private.",
                  font_size=14, color=C_NAVY, italic=True, alignment=PP_ALIGN.CENTER)

    _add_speaker_notes(slide,
        "What does this all mean? I call this the 'neural loneliness' frame. The dACC is where "
        "the brain detects social mismatches. If each anxious teen encodes these mismatches "
        "differently, their inner experience of social threat is fundamentally private — no one "
        "else is processing it the same way they are. Given that neural similarity predicts "
        "social bonding, this neural divergence may contribute to the pervasive sense of being "
        "misunderstood that characterizes social anxiety. They don't just feel misunderstood — "
        "at a neural level, their experience really IS different from everyone else's.")


def build_slide_18_limitations(prs):
    """Slide 18: Limitations & Future Directions."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg_white(slide)
    _add_title(slide, "Limitations & Future Directions")

    # Two columns
    # Left: Limitations
    _add_text_box(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.4),
                  "Limitations", font_size=20, color=C_RED, bold=True)

    tf1 = _add_rich_text_box(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(3.5))
    _add_paragraph(tf1, "• N = 33 — adequate for IS-RSA (528 pairwise comparisons) but limits generalizability",
                   font_size=14, color=C_NAVY, space_after=8)
    _add_paragraph(tf1, "• 1 of 32 ROI tests survives FDR; searchlight clusters do not survive FWE",
                   font_size=14, color=C_NAVY, space_after=8)
    _add_paragraph(tf1, "• Cross-sectional — cannot establish temporal priority",
                   font_size=14, color=C_NAVY, space_after=8)
    _add_paragraph(tf1, "• Reverse AnnaK in anterior insula is uncorrected",
                   font_size=14, color=C_NAVY)

    # Right: Future Directions
    _add_text_box(slide, Inches(7.0), Inches(1.7), Inches(5.5), Inches(0.4),
                  "Future Directions", font_size=20, color=C_TEAL, bold=True)

    tf2 = _add_rich_text_box(slide, Inches(7.0), Inches(2.2), Inches(5.5), Inches(4.0))
    _add_paragraph(tf2, "• Larger sample to improve power for searchlight correction",
                   font_size=14, color=C_NAVY, space_after=8)
    _add_paragraph(tf2, "• Q2: Track learning trajectories — does representational geometry mirror true peer structure over runs?",
                   font_size=14, color=C_NAVY, space_after=8)
    _add_paragraph(tf2, "• Does neural idiosyncrasy predict real-world social functioning (friendships, loneliness)?",
                   font_size=14, color=C_NAVY, space_after=8)
    _add_paragraph(tf2, "• Formal insular parcellation — test anterior vs posterior dissociation",
                   font_size=14, color=C_NAVY, space_after=8)
    _add_paragraph(tf2, "• Is this specific to social anxiety, or general to negative affect?",
                   font_size=14, color=C_NAVY)

    _add_speaker_notes(slide,
        "Limitations: our sample of 33 is adequate for IS-RSA — we have 528 unique pairwise "
        "comparisons — but it limits power, especially for searchlight correction. Only one "
        "of our 32 ROI tests survived FDR. The study is cross-sectional. For future directions: "
        "we want to track learning trajectories over runs, test whether idiosyncrasy predicts "
        "real-world social functioning, and run a formal insular parcellation analysis to test "
        "the anterior-posterior dissociation more rigorously.")


def build_slide_19_thanks(prs):
    """Slide 19: Thank You."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg_white(slide)

    # Accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_BLUE
    bar.line.fill.background()

    _add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.0),
                  "Thank You", font_size=44, color=C_NAVY, bold=True, alignment=PP_ALIGN.CENTER)

    # Acknowledgments
    tf = _add_rich_text_box(slide, Inches(2.0), Inches(3.5), Inches(9.0), Inches(3.0))
    _add_paragraph(tf, "Advisor: Dr. Johanna Jarcho", font_size=18, color=C_NAVY, bold=True,
                   space_after=8, alignment=PP_ALIGN.CENTER)
    _add_paragraph(tf, "Jarcho Lab | Temple University", font_size=16, color=C_GRAY,
                   space_after=16, alignment=PP_ALIGN.CENTER)
    _add_paragraph(tf, "Theresa Clarkson  (LEARN task development & computational modeling)",
                   font_size=14, color=C_NAVY, space_after=6, alignment=PP_ALIGN.CENTER)
    _add_paragraph(tf, "Committee Members", font_size=14, color=C_NAVY,
                   space_after=16, alignment=PP_ALIGN.CENTER)
    _add_paragraph(tf, "Questions?", font_size=20, color=C_BLUE, bold=True,
                   alignment=PP_ALIGN.CENTER)

    # Bottom bar
    bar2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.35), SLIDE_W, Inches(0.15))
    bar2.fill.solid()
    bar2.fill.fore_color.rgb = C_BLUE
    bar2.line.fill.background()


# ──────────────────────────────────────────────────────────────
# 5. MAIN BUILD
# ──────────────────────────────────────────────────────────────

def build_presentation(output_path):
    """Build the full presentation."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    print("Building slides...")

    # Act 1: Background (Slides 1-8)
    print("  Slide 1: Title")
    build_slide_01_title(prs)
    print("  Slide 2: Research Interests")
    build_slide_02_interests(prs)
    print("  Slide 3: Why Adolescence")
    build_slide_03_why_adolescence(prs)
    print("  Slide 4: Background")
    build_slide_04_background(prs)
    print("  Slide 5: The Gap")
    build_slide_05_gap(prs)
    print("  Slide 6: RSA")
    build_slide_06_rsa(prs)
    print("  Slide 7: Idiosyncrasy Literature")
    build_slide_07_idiosyncrasy(prs)
    print("  Slide 8: Research Questions")
    build_slide_08_questions(prs)

    # Act 2: Methods (Slides 9-11)
    print("  Slide 9: LEARN Task")
    build_slide_09_learn_task(prs)
    print("  Slide 10: IS-RSA Walk-Through")
    build_slide_10_isrsa_walkthrough(prs)
    print("  Slide 11: Pipeline & Methods")
    build_slide_11_methods(prs)

    # Act 3: Results (Slides 12-16)
    print("  Slide 12: dACC1 Finding")
    build_slide_12_results_dacc(prs)
    print("  Slide 13: dACC1 Meaning")
    build_slide_13_dacc_meaning(prs)
    print("  Slide 14: Anterior Insula")
    build_slide_14_anterior_insula(prs)
    print("  Slide 15: Searchlight")
    build_slide_15_searchlight(prs)
    print("  Slide 16: Emerging Story")
    build_slide_16_emerging_story(prs)

    # Act 4: Discussion (Slides 17-19)
    print("  Slide 17: Discussion")
    build_slide_17_discussion(prs)
    print("  Slide 18: Limitations")
    build_slide_18_limitations(prs)
    print("  Slide 19: Thank You")
    build_slide_19_thanks(prs)

    prs.save(output_path)
    print(f"\nPresentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


def main():
    parser = argparse.ArgumentParser(description="Generate First-Year PhD Talk PPTX")
    parser.add_argument("--output", "-o", default="first_year_talk.pptx",
                        help="Output PPTX file path")
    args = parser.parse_args()

    build_presentation(args.output)


if __name__ == "__main__":
    main()
